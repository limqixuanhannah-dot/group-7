#!/usr/bin/env python3
"""Build a 7-slide deck: SENTRY Phase 2 — New Capabilities, Monitoring & Intelligence Update."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(0x0A, 0x0E, 0x17)
DARK_CARD = RGBColor(0x11, 0x18, 0x27)
BORDER = RGBColor(0x1E, 0x29, 0x3B)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
TEXT_PRIMARY = RGBColor(0xE0, 0xE6, 0xF0)
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=DARK_CARD, line_color=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def slide_number(slide, num):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4),
                 f"{num:02d}", font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)


# ─── SLIDE 1: TITLE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
             "◆  SENTRY", font_size=16, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(4), Inches(0.4),
             "PHASE 2 — CAPABILITY EXPANSION & RESEARCH UPDATE", font_size=10, color=TEXT_MUTED)

add_text_box(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(1.2),
             "SENTRY Phase 2:\nNew Capabilities & Threat Landscape",
             font_size=46, bold=True, color=WHITE)

tf = add_text_box(slide, Inches(0.8), Inches(3.8), Inches(10), Inches(0.8),
                  "Project monitoring, intelligence research, and 12 new operational systems\nadded between 16–19 June 2026.",
                  font_size=18, color=TEXT_SECONDARY)

add_text_box(slide, Inches(0.8), Inches(4.8), Inches(6), Inches(0.4),
             "LIVE   |   DSTA-FOCUSED   |   NEW CAPABILITIES DEPLOYED",
             font_size=11, bold=True, color=ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(5), Inches(0.4),
             "Group 7  ·  Hannah (Principal Operator)  ·  19 June 2026",
             font_size=11, color=TEXT_MUTED)
slide_number(slide, 1)


# ─── SLIDE 2: HOURLY MONITORING & BREACH INTELLIGENCE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
             "MONITORING & INTELLIGENCE", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Hourly Breach Monitoring Pipeline",
             font_size=34, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.5),
             "Automated hourly scans deliver real-time cyber threat intelligence to Discord #group-7.",
             font_size=15, color=TEXT_SECONDARY)

# Pipeline flow cards
pipeline = [
    ("⏰  Every Hour", "Cron-triggered scan\nagainst configured\ntarget domains & emails"),
    ("🔍  Breach Check", "HIBP API lookup\n+ web intelligence\nfor new exposures"),
    ("📰  Intel Roundup", "Cyber security news\nsourced from:\nNVD, CISA, BleepingComputer"),
    ("📡  Delivery", "Formatted report to\nDiscord #group-7\nwith severity flags")
]

for i, (title, desc) in enumerate(pipeline):
    left = Inches(0.8 + i * 3.1)
    card = add_shape(slide, left, Inches(2.6), Inches(2.9), Inches(2.0))
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.font.name = 'Calibri'
    p2.space_before = Pt(8)

# Breach highlights
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(6), Inches(0.4),
             "KEY BREACHES DETECTED (16–19 JUNE 2026)",
             font_size=11, bold=True, color=ACCENT_BLUE)

breaches = [
    ("Canvas/Instructure", "275M users — ShinyHunters — 3.65 TB data", ACCENT_RED),
    ("Kodak", "2.2M records — ShinyHunters ransomware claimed", ACCENT_RED),
    ("Carnival Corp", "5.99M guests — unauthorized employee account", ACCENT_ORANGE),
    ("ADT Security", "5.5M records leaked — customer contact data", ACCENT_ORANGE),
    ("ShinyHunters Vishing", "100+ orgs targeted — MFA bypass via voice phishing", ACCENT_ORANGE),
]

for i, (name, detail, color) in enumerate(breaches):
    top = Inches(5.4 + i * 0.38)
    # Name tag
    tag = add_shape(slide, Inches(0.8), top, Inches(2.8), Inches(0.32),
                    fill_color=DARK_CARD, line_color=color)
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    # Detail
    add_text_box(slide, Inches(3.8), top, Inches(8), Inches(0.32),
                 detail, font_size=10, color=TEXT_SECONDARY)

slide_number(slide, 2)


# ─── SLIDE 3: PROJECT RESEARCH ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
             "RESEARCH & INTELLIGENCE", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Active Project Research",
             font_size=34, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.5),
             "Hourly research cycles correlate web intelligence to active projects and append findings.",
             font_size=15, color=TEXT_SECONDARY)

# Project Fish card
fish_card = add_shape(slide, Inches(0.8), Inches(2.4), Inches(5.8), Inches(3.8))
tf = fish_card.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🐟  Project: Fish (owner: hannah)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN
p.font.name = 'Calibri'

fish_items = [
    "Practical Fishkeeping June 2026 issue articles logged",
    "Blue Space Effect — tanks lower BP in 5 min (UK research)",
    "Nano Tank Ecosystems — 5-10 gal management guidance",
    "CARES Conservation — 30+ fish species saved by aquarists",
    "Skiffia francesae successfully reintroduced from extinction",
]
for item in fish_items:
    p2 = tf.add_paragraph()
    p2.text = f"→  {item}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.font.name = 'Calibri'
    p2.space_before = Pt(4)

# Project Google Slides card
gs_card = add_shape(slide, Inches(7.0), Inches(2.4), Inches(5.8), Inches(3.8))
tf = gs_card.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📊  Project: Google Slides (owner: hannah, collab: nc)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_ORANGE
p.font.name = 'Calibri'

gs_items = [
    "Gemini Classroom app integration — GA (17 Jun)",
    "Gmail as Ask Gemini source — GA (3 Jun)",
    "Calendar DLP policies — GA — scan events for sensitive data",
    "New Gmail Security Advisor for SMBs",
    "Gemini Workspace outage (10 Jun) — 1,300+ reports",
    "Workspace Policy API — DLP Mutate Endpoints (8 Jun)",
    "ShinyHunters vishing targeting Workspace/M365/Okta",
]
for item in gs_items:
    p2 = tf.add_paragraph()
    p2.text = f"→  {item}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.font.name = 'Calibri'
    p2.space_before = Pt(4)

# Key threats bar
threats_box = add_shape(slide, Inches(0.8), Inches(6.5), Inches(12.0), Inches(0.6),
                        fill_color=RGBColor(0x1a, 0x0a, 0x0a), line_color=ACCENT_RED)
tf = threats_box.text_frame
p = tf.paragraphs[0]
p.text = "⚠️  Supply-chain attacks escalating: Miasma worm leaked · Mastra AI npm (1.1M downloads) · JetBrains plugins stealing API keys · Open source registries under fire"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = ACCENT_RED
p.font.name = 'Calibri'

slide_number(slide, 3)


# ─── SLIDE 4: 12 NEW SYSTEMS OVERVIEW ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
             "OPERATIONAL EXPANSION", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "12 New Operational Systems",
             font_size=34, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.5),
             "All systems implemented and deployed on 19 June 2026 per Group 7 request.",
             font_size=15, color=TEXT_SECONDARY)

systems = [
    ("SECURITY & COMPLIANCE", [
        ("🛡️", "Classification Tagging", "RESTRICTED/CONFIDENTIAL/SECRET per project"),
        ("📅", "Compliance Deadlines", "Auto-detect overdue/due-soon/upcoming"),
        ("🏢", "Vendor Posture Tracker", "Clearance expiry, assessment, incident log"),
    ], ACCENT_CYAN),
    ("OPS & INCIDENT RESPONSE", [
        ("🚨", "IR Drill Scheduling", "Schedule, score, track participants"),
        ("📋", "Crisis Comms Templates", "Breach/Outage/Disruption — fill & send"),
    ], ACCENT_ORANGE),
    ("PROJECT INTELLIGENCE", [
        ("📊", "Risk Scoring Dashboard", "4-dimension composite risk bands"),
        ("🌐", "Threat Briefings", "Correlate active threats to projects"),
        ("🐛", "Vulnerability Triage", "PENDING → IN_PROGRESS → PATCHED"),
    ], ACCENT_RED),
    ("LOGISTICS & GOVERNANCE", [
        ("💻", "Asset Registry", "Hardware/software with ownership & expiry"),
        ("📝", "Secure File Transfer Log", "CREATE/ACCESS/SHARE/EDIT/DELETE logs"),
        ("🎯", "Action Item Tracker", "Extract from meeting notes, auto-date"),
        ("📜", "Decision Log", "Timestamped audit trail with rationale"),
    ], ACCENT_GREEN),
]

y = 2.2
for section, items, color in systems:
    # Section label
    add_text_box(slide, Inches(0.8), Inches(y), Inches(3.0), Inches(0.3),
                 section, font_size=9, bold=True, color=color)
    x = 0.8
    for icon, title, desc in items:
        card = add_shape(slide, Inches(x), Inches(y + 0.3), Inches(3.0), Inches(1.1),
                        fill_color=DARK_CARD, line_color=BORDER)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_PRIMARY
        p.font.name = 'Calibri'
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9)
        p2.font.color.rgb = TEXT_MUTED
        p2.font.name = 'Calibri'
        p2.space_before = Pt(2)
        x += 3.15
    y += 1.6

slide_number(slide, 4)


# ─── SLIDE 5: KEY THREAT INTELLIGENCE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
             "THREAT LANDSCAPE", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Cyber Threat Intelligence — June 2026",
             font_size=34, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.4),
             "Major findings from hourly monitoring cycles (16–19 June).",
             font_size=14, color=TEXT_SECONDARY)

threats = [
    ("SUPPLY CHAIN ATTACKS", ACCENT_RED, [
        "Miasma worm leaked via compromised GitHub — targets PyPI, npm, RubyGems, GitHub Actions",
        "Mastra AI npm package (1.1M weekly downloads) — malicious code in AI tooling",
        "JetBrains Marketplace — 15 plugins stealing DeepSeek/OpenAI API keys",
        "Security tools compromised: Trivy, Bitwarden CLI, Checkmarx → downstream OpenAI/Vercel breaches",
    ]),
    ("ACTIVELY EXPLOITED VULNERABILITIES", ACCENT_ORANGE, [
        "FortiSandbox — 3 critical flaws under active exploitation (CISA KEV)",
        "Cisco SD-WAN Manager CVE-2026-20262 (CVSS 6.5) — patches released",
        "Oracle PeopleSoft CVE-2026-35273 — zero-day actively exploited in wild",
        "Android zero-day — June 2026 patch cycle addressed active exploitation",
    ]),
    ("EMERGING THREATS & TRENDS", ACCENT_CYAN, [
        "ShinyHunters: 14 of 37 major breaches in 2026 — vishing + social engineering primary vector",
        "HTTP/2 Bomb (CVE-2026-45247) — NGINX, Apache, IIS, Envoy, Cloudflare Pingora affected",
        "Rokarolla Android trojan — 217 banking/crypto apps targeted via 137 commands",
        "PCPJack worm — credential-stealing targeting exposed Docker, K8s, Redis, MongoDB",
    ]),
]

for i, (section, color, items) in enumerate(threats):
    top = Inches(2.2 + i * 1.75)
    # Section header
    add_text_box(slide, Inches(0.8), top, Inches(4), Inches(0.3),
                 section, font_size=10, bold=True, color=color)
    # Items
    for j, item in enumerate(items):
        add_text_box(slide, Inches(1.2), Inches(top.inches + j * 0.4 + 0.35), Inches(11.5), Inches(0.35),
                     f"→  {item}", font_size=10, color=TEXT_SECONDARY)

slide_number(slide, 5)


# ─── SLIDE 6: SENTRY FRAMEWORK EXPANSION ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
             "FRAMEWORK EVOLUTION", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "SENTRY Framework — Expanded",
             font_size=34, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.5),
             "From 5 to 9 capability domains. All new systems are CLI-accessible and auditable.",
             font_size=15, color=TEXT_SECONDARY)

# Original 5 SENTRY domains
orig = [
    ("🛡️  Cyber Hygiene\n& Breach Monitoring", "Phase 1"),
    ("🚨  Crisis Triage", "Phase 1"),
    ("🔎  Misinformation\n& Fact Check", "Phase 1"),
    ("🌐  OSINT & Geopolitical\nAnalysis", "Phase 1"),
    ("⚡  Ops Readiness\n& Task Tracking", "Phase 1"),
]

# New 4 domains
new = [
    ("📜  Security &\nCompliance", "Phase 2"),
    ("🚨  IR Drills &\nCrisis Templates", "Phase 2"),
    ("🧠  Project Intelligence\nRisk & Vulnerability", "Phase 2"),
    ("📦  Asset & Governance\nLogistics Tracking", "Phase 2"),
]

# Original 5
for i, (title, phase) in enumerate(orig):
    left = Inches(0.8 + i * 2.5)
    card = add_shape(slide, left, Inches(2.4), Inches(2.3), Inches(1.8))
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = phase
    p2.font.size = Pt(9)
    p2.font.color.rgb = ACCENT_BLUE
    p2.font.name = 'Calibri'
    p2.space_before = Pt(12)

# Arrow
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(12), Inches(0.4),
             "⬇  ORIGINAL 5 DOMAINS  ⬇         ⬇  +4 NEW DOMAINS  ⬇",
             font_size=12, bold=True, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# New 4
for i, (title, phase) in enumerate(new):
    left = Inches(1.8 + i * 2.8)
    card = add_shape(slide, left, Inches(5.0), Inches(2.6), Inches(1.8),
                    fill_color=RGBColor(0x14, 0x1D, 0x2F), line_color=ACCENT_CYAN)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = phase
    p2.font.size = Pt(9)
    p2.font.color.rgb = ACCENT_CYAN
    p2.font.name = 'Calibri'
    p2.space_before = Pt(12)

slide_number(slide, 6)


# ─── SLIDE 7: NEXT STEPS ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
             "NEXT STEPS", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "ROADMAP — Activate & Integrate",
             font_size=34, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.4),
             "Immediate priorities to move from deployed to operational.",
             font_size=14, color=TEXT_SECONDARY)

# Three priority boxes
priorities = [
    ("CRITICAL", ACCENT_RED, "Configure Monitoring Targets",
     "Add domains / emails to monitoring/config.json\nEnable HIBP API key\nActivate breach detection pipeline\n\nWithout targets, the hourly scanner runs\nidle — 48+ reports sent with zero findings."),
    ("HIGH", ACCENT_ORANGE, "Integrate Spreadsheet",
     "Link Google Sheet for task/deliverable tracking\nEnable Claw read/write access\nAuto-update status from sheet changes\n\nFoundation for the 12 new systems to\nread/write project data."),
    ("HIGH", ACCENT_ORANGE, "Renew Discord Token",
     "Generate new bot token (401 error)\nRestore Discord delivery pipeline\nEnable automated report delivery\nto #group-7 channel."),
]

for i, (severity, color, title, desc) in enumerate(priorities):
    left = Inches(0.8 + i * 4.1)
    card = add_shape(slide, left, Inches(2.4), Inches(3.8), Inches(4.5))
    # Severity tag
    tag = add_shape(slide, Inches(left.inches + 0.3), Inches(2.7), Inches(1.6), Inches(0.3),
                    fill_color=RGBColor(0x2a, 0x0a, 0x0a) if severity == "CRITICAL" else RGBColor(0x2a, 0x1a, 0x0a),
                    line_color=color)
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = severity
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, Inches(left.inches + 0.3), Inches(3.2), Inches(3.2), Inches(0.4),
                 title, font_size=15, bold=True, color=WHITE)

    # Description
    add_text_box(slide, Inches(left.inches + 0.3), Inches(3.7), Inches(3.2), Inches(2.8),
                 desc, font_size=11, color=TEXT_SECONDARY)

# Bottom banner
banner = add_shape(slide, Inches(0.8), Inches(6.8), Inches(12.0), Inches(0.5),
                   fill_color=RGBColor(0x0F, 0x17, 0x2A), line_color=ACCENT_BLUE)
tf = banner.text_frame
p = tf.paragraphs[0]
p.text = "12 new systems deployed · 9 SENTRY domains live · 5+ breach sources monitored · 2 active projects researched hourly"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER

slide_number(slide, 7)


# ─── SAVE ───
output_path = os.path.expanduser('~/.openclaw/workspace/SENTRY-Phase2-Update-Deck.pptx')
prs.save(output_path)
print(f'Slide count: {len(prs.slides)}')
for i, slide in enumerate(prs.slides):
    shapes = len(slide.shapes)
    print(f'  Slide {i+1}: {shapes} shapes')
print(f'Saved to: {output_path}')
