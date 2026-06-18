#!/usr/bin/env python3
"""Build Project Claw 10-slide PPTX from the project-claw-deck.html content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(0x0A, 0x0E, 0x17)
DARK_CARD = RGBColor(0x11, 0x18, 0x27)
BORDER = RGBColor(0x1E, 0x29, 0x3B)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
TEXT_PRIMARY = RGBColor(0xE0, 0xE6, 0xF0)
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=DARK_CARD, line_color=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=14, bold=False, color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT, space_before=4, space_after=4, font_name='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def slide_number(slide, num):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4),
                 f"{num:02d}", font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)


# ─── SLIDE 1: TITLE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Logo marker
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
             "◆  PROJECT CLAW", font_size=16, bold=True, color=ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(4), Inches(0.4),
             "CRISIS INFORMATION TRIAGE & OPS-READINESS", font_size=10, color=TEXT_MUTED)

# Title
add_text_box(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(1.2),
             "A Claw for Defence & Security\nDecision Support",
             font_size=48, bold=True, color=WHITE)

# Subtitle
tf = add_text_box(slide, Inches(0.8), Inches(3.6), Inches(10), Inches(0.8),
                  "An AI-powered project management assistant that helps teams\ndecide faster, act smarter, and stay operationally ready.",
                  font_size=18, color=TEXT_SECONDARY)

# Tags
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(6), Inches(0.4),
             "LIVE PROTOTYPE   |   DSTA-FOCUSED   |   SECURITY-GRADE   |   SINGAPORE-CENTRIC",
             font_size=11, bold=True, color=ACCENT_BLUE)

# Footer
add_text_box(slide, Inches(0.8), Inches(6.8), Inches(5), Inches(0.4),
             "Group 7  ·  Hannah (Principal Operator)  ·  June 2026",
             font_size=11, color=TEXT_MUTED)
add_text_box(slide, Inches(0.8), Inches(7.0), Inches(4), Inches(0.3),
             "CONFIDENTIAL — For DSTA Project Review",
             font_size=9, color=RGBColor(0x33, 0x41, 0x55))
slide_number(slide, 1)


# ─── SLIDE 2: PROBLEM STATEMENT ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
             "PROBLEM STATEMENT", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
             "The Information Overload Problem",
             font_size=36, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(10), Inches(0.6),
             "Large projects generate massive information flow — chat messages, task updates, pinned items, status changes, purchase requests. Teams and executive committees struggle to maintain situational awareness.",
             font_size=16, color=TEXT_SECONDARY)

# Problem box
box = add_shape(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.0),
                fill_color=RGBColor(0x0F, 0x17, 0x2A), line_color=ACCENT_BLUE)
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"How can we create an interactive project manager for members and executive committees of large projects that can make updates, new information, and tasks easier to overview?"'
p.font.size = Pt(16)
p.font.color.rgb = TEXT_PRIMARY
p.font.italic = True
p.font.name = 'Calibri'

# Three cards
cards_data = [
    ("📊  Information Fragmentation", "Updates spread across chat, email, spreadsheets, documents. No single pane of glass."),
    ("⏱️  Decision Latency", "Critical information gets buried. Teams waste time searching instead of acting."),
    ("🔍  Priority Blindness", "Urgent vs important gets blurred. No automated triage or escalation system.")
]

for i, (title, desc) in enumerate(cards_data):
    left = Inches(0.8 + i * 4.0)
    card = add_shape(slide, left, Inches(3.9), Inches(3.7), Inches(1.8))
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.font.name = 'Calibri'
    p2.space_before = Pt(8)

slide_number(slide, 2)


# ─── SLIDE 3: SOLUTION OVERVIEW ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
             "SOLUTION OVERVIEW", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "What Is The Claw?",
             font_size=36, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(10), Inches(0.6),
             "An AI-powered ops assistant that ingests project chatter, triages information, tracks tasks, checks logistics, and surfaces what matters — in real time.",
             font_size=16, color=TEXT_SECONDARY)

# 4 cards in a row
cards = [
    ("✂️  Summarise & Filter", "Compress noisy information streams into clear, actionable briefs. Noise goes out."),
    ("🚩  Flag Priorities", "Identify urgency, priority, and risk. CRITICAL — HIGH — MEDIUM — LOW."),
    ("📋  Track Tasks", "Operational awareness. PENDING — IN PROGRESS — COMPLETE — BLOCKED."),
    ("✅  Ops Readiness", "Checklists, drills, preparedness scoring, gap analysis. Know your status at all times.")
]

for i, (title, desc) in enumerate(cards):
    left = Inches(0.8 + i * 3.1)
    card = add_shape(slide, left, Inches(2.8), Inches(2.9), Inches(2.2))
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.font.name = 'Calibri'
    p2.space_before = Pt(8)

# Status bar items
metrics = [
    ("5", "Core Capability Domains"),
    ("24/7", "Operational Availability"),
    ("4", "Priority Tiers"),
    ("4", "Task Status States")
]
for i, (num, label) in enumerate(metrics):
    left = Inches(1.5 + i * 2.8)
    box = add_shape(slide, left, Inches(5.5), Inches(2.4), Inches(1.2),
                    fill_color=DARK_CARD, line_color=BORDER)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = 'Calibri'
    p2.alignment = PP_ALIGN.CENTER

slide_number(slide, 3)


# ─── SLIDE 4: CORE CAPABILITIES (SENTRY) ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
             "CORE CAPABILITIES", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Full SENTRY Framework",
             font_size=36, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(10), Inches(0.4),
             "Five integrated defence and security capabilities — all operational.",
             font_size=14, color=TEXT_SECONDARY)

sentry_cards = [
    ("🛡️  Cyber Hygiene & Breach", "Password audit, 2FA checks, account exposure scans. HIBP integration for real-time breach alerts."),
    ("🚨  Crisis Triage", "Step-by-step incident playbooks, resource mapping, decision matrices. KNOW — DECIDE — ACT."),
    ("🔎  Misinformation & Fact Check", "Cross-reference claims against Tier 1-4 sources. Source credibility scoring. Deepfake analysis."),
    ("🌐  OSINT & Geopolitical Analysis", "Situation briefs, threat actor profiles, capability assessments. Tiered sourcing from official to open."),
    ("⚡  Ops Readiness & Task Tracking", "Preparedness scoring, checklists. Priority-triage project management. Spreadsheet integration."),
    ("🗂️  Spreadsheet & Logistics", "Reads/writes project spreadsheets. Tracks missing items, purchase status, inventory gaps.")
]

for i, (title, desc) in enumerate(sentry_cards):
    row = i // 3
    col = i % 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(2.4 + row * 2.3)
    card = add_shape(slide, left, top, Inches(3.8), Inches(2.0))
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.font.name = 'Calibri'
    p2.space_before = Pt(8)

slide_number(slide, 4)


# ─── SLIDE 5: PROJECT STATUS ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
             "PROJECT STATUS", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Current Operational Status",
             font_size=36, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(10), Inches(0.4),
             "Live status as of June 2026, UTC+8.",
             font_size=14, color=TEXT_MUTED)

# Status items
statuses = [
    ("✅", "Dashboard", "Operational"),
    ("✅", "Heartbeat System", "Online"),
    ("⚠️", "Discord Token", "Needs Renewal"),
    ("✅", "SENTRY Framework", "Deployed"),
    ("✅", "Breach Monitoring", "Ready"),
    ("⚠️", "Excel Tracking", "Needs Link")
]

for i, (icon, name, status) in enumerate(statuses):
    row = i // 3
    col = i % 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(2.3 + row * 1.3)
    box = add_shape(slide, left, top, Inches(3.8), Inches(1.0),
                    fill_color=DARK_CARD, line_color=BORDER)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{icon}  {name}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = f"Status: {status}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = GREEN if status == "Online" or status == "Operational" or status == "Deployed" or status == "Ready" else YELLOW
    p2.font.name = 'Calibri'

# To-Do table header
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(4), Inches(0.4),
             "PRIORITY & URGENCY — SORTED TO-DO LIST",
             font_size=11, bold=True, color=ACCENT_BLUE)

todo_data = [
    ("CRITICAL", "Renew Discord bot token (401 Unauthorized)", "PENDING", "Hannah"),
    ("CRITICAL", "Restore memory index (embedding mismatch)", "PENDING", "Hannah"),
    ("HIGH", "Set up Excel/Sheets project tracking", "IN PROGRESS", "Group 7"),
    ("HIGH", "Confirm logistics purchases", "PENDING", "Hannah"),
    ("HIGH", "Create presentation deck (this deck ✓)", "COMPLETE", "Group 7"),
    ("MEDIUM", "Compile DSTA-aligned use cases", "COMPLETE", "Group 7"),
    ("MEDIUM", "Build web dashboard frontend", "PENDING", "Deferred"),
    ("LOW", "Stale dashboard entry — cleanup", "PENDING", "Housekeeping")
]

# Draw table manually with shapes
headers = ["Priority", "Task", "Status", "Owner"]
cols_w = [1.6, 6.0, 1.6, 1.6]
x_start = 0.8
y_start = 5.4

# Header row
x = x_start
for j, (h, w) in enumerate(zip(headers, cols_w)):
    box = add_shape(slide, Inches(x), Inches(y_start), Inches(w), Inches(0.35),
                    fill_color=RGBColor(0x14, 0x1D, 0x2F), line_color=BORDER)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = h
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = TEXT_MUTED
    p.font.name = 'Calibri'
    x += w

for i, row_data in enumerate(todo_data):
    y = y_start + 0.35 + i * 0.38
    x = x_start
    colors = {
        "CRITICAL": ACCENT_RED,
        "HIGH": ACCENT_ORANGE,
        "MEDIUM": ACCENT_BLUE,
        "LOW": TEXT_MUTED
    }
    for j, (val, w) in enumerate(zip(row_data, cols_w)):
        box = add_shape(slide, Inches(x), Inches(y), Inches(w), Inches(0.38),
                        fill_color=DARK_CARD, line_color=RGBColor(0x1A, 0x23, 0x40))
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(10)
        clr = colors.get(val, TEXT_PRIMARY) if j == 0 else TEXT_PRIMARY
        if j == 3:
            clr = TEXT_SECONDARY
        p.font.color.rgb = TEXT_PRIMARY if j > 0 else clr
        p.font.bold = (j == 0)
        p.font.name = 'Calibri'
        x += w

slide_number(slide, 5)


# ─── SLIDE 6: CHAT SUMMARY ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
             "INTELLIGENCE SUMMARY", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Chat Summary & Pinned Highlights",
             font_size=36, bold=True, color=WHITE)

# Pinned items
pinned = [
    ("SUPER IMPORTANT — create a professional presentation on building a claw for defence/security", "Hannah flagged this as top priority. Deck built and served."),
    ("Link to excel spreadsheet to track missing items — Claw can update when prompted", "Spreadsheet system needs linking. CSV template ready."),
    ("Check logistics — check if you bought stuff", "Logistics audit requested. Pending confirmation of orders vs deliveries.")
]

add_text_box(slide, Inches(0.8), Inches(1.7), Inches(4), Inches(0.4),
             "📌  PINNED HIGHLIGHTS",
             font_size=11, bold=True, color=ACCENT_RED)

for i, (title, desc) in enumerate(pinned):
    left = Inches(0.8)
    top = Inches(2.2 + i * 1.3)
    box = add_shape(slide, left, top, Inches(11.5), Inches(1.1),
                    fill_color=DARK_CARD, line_color=RGBColor(0xEF, 0x44, 0x44))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"📌 {title}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.font.name = 'Calibri'
    p2.space_before = Pt(4)

# Session summary
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(4), Inches(0.4),
             "📋  SESSION KEY POINTS",
             font_size=11, bold=True, color=ACCENT_BLUE)

summary_items = [
    "Create professional presentation on building a Claw for defence/security",
    "Update on urgent project items & sort to-do list by priority",
    "Link spreadsheet to track missing items — Claw updates on demand",
    "Check logistics — what was ordered vs received",
    "Target audience: DSTA project people / executive committees"
]

for i, item in enumerate(summary_items):
    add_text_box(slide, Inches(1.2), Inches(6.1 + i * 0.28), Inches(10), Inches(0.3),
                 f"→  {item}", font_size=11, color=TEXT_SECONDARY)

slide_number(slide, 6)


# ─── SLIDE 7: LOGISTICS AUDIT ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
             "LOGISTICS AUDIT", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Purchase & Inventory Status",
             font_size=36, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(10), Inches(0.6),
             "No existing purchase records found. Below is the tracking system — The Claw can update when prompted with what was ordered, received, or is missing.",
             font_size=14, color=TEXT_SECONDARY)

# Two cards
left_card = add_shape(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(2.5))
tf = left_card.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📊  Project Tracker Spreadsheet"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = TEXT_PRIMARY
p.font.name = 'Calibri'
for line in [
    "File: project-claw-tracker.csv",
    "Format: CSV (import to Excel / Google Sheets)",
    "Auto-update: The Claw reads/writes on demand",
    "",
    "Ready for data input — tell me what was ordered"
]:
    p2 = tf.add_paragraph()
    p2.text = line
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.font.name = 'Calibri'
    p2.space_before = Pt(4)

right_card = add_shape(slide, Inches(6.8), Inches(2.6), Inches(5.5), Inches(2.5))
tf = right_card.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📦  Items to Track"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = TEXT_PRIMARY
p.font.name = 'Calibri'
for line in [
    "No purchase data found yet.",
    "",
    "Tell me what was ordered and I'll populate",
    "the tracker. I can check status, flag",
    "missing items, and alert on delivery gaps."
]:
    p2 = tf.add_paragraph()
    p2.text = line
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.font.name = 'Calibri'
    p2.space_before = Pt(4)

# Callout box
box = add_shape(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8),
                fill_color=RGBColor(0x0F, 0x17, 0x2A), line_color=ACCENT_BLUE)
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "To start tracking: tell me \"I bought X from Y on Z date\" or \"Check if item A was delivered\" and I'll update the spreadsheet. I can also scan for missing items automatically."
p.font.size = Pt(13)
p.font.color.rgb = TEXT_PRIMARY
p.font.name = 'Calibri'

slide_number(slide, 7)


# ─── SLIDE 8: TECHNICAL ARCHITECTURE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
             "TECHNICAL ARCHITECTURE", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "How The Claw Works",
             font_size=36, bold=True, color=WHITE)

arch_cards = [
    ("🧠  AI Core (OpenClaw)", "DeepSeek V4 reasoning engine. Sub-agent orchestration. Memory system for session continuity. Tool-calling interface."),
    ("🔧  Tool Layer", "File read/write for spreadsheets. Web search & fetch for real-time intel. Cron scheduling. Sessions for parallel delegation."),
    ("📡  Communication Layer", "WebChat direct. Discord (pending token). Web dashboard. Multi-surface: direct, group, isolated sessions."),
    ("🗄️  Data Layer", "File-based workspace storage. Daily memory + curated MEMORY.md. CSV tracking. SENTRY framework doctrine (project67.md).")
]

for i, (title, desc) in enumerate(arch_cards):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(2.2 + row * 2.2)
    card = add_shape(slide, left, top, Inches(5.8), Inches(1.9))
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.font.name = 'Calibri'
    p2.space_before = Pt(8)

# Status checks
tech_status = [
    ("✓", "Multi-Session\nOrchestration"),
    ("✓", "File I/O for\nSpreadsheets"),
    ("✓", "Web Search\n& Fetch"),
    ("⚠", "Discord\nIntegration"),
    ("✓", "Cron / Scheduled\nTasks")
]

for i, (icon, label) in enumerate(tech_status):
    left = Inches(0.8 + i * 2.5)
    box = add_shape(slide, left, Inches(6.2), Inches(2.2), Inches(0.9),
                    fill_color=DARK_CARD, line_color=BORDER)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{icon}  {label}"
    p.font.size = Pt(10)
    p.font.color.rgb = GREEN if icon == "✓" else YELLOW
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

slide_number(slide, 8)


# ─── SLIDE 9: ROADMAP ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
             "DEVELOPMENT ROADMAP", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "From Prototype to Production",
             font_size=36, bold=True, color=WHITE)

roadmap = [
    ("PHASE 1  ✓", ACCENT_BLUE, "Core Foundation", "SENTRY framework deployed. Identity, personality, memory configured. Crisis triage + OSINT + cyber hygiene operational. Heartbeat system online."),
    ("PHASE 1  ✓", ACCENT_BLUE, "Breach Monitoring Module", "HIBP integration designed. Watchlist system, alert pipeline, priority matrix documented. Ready for API key activation."),
    ("PHASE 2  ← NOW", ACCENT_ORANGE, "Project Management Layer", "Spreadsheet integration for tracking. To-do system with priority triage. Logistics/purchase tracking. Presentation deck for DSTA."),
    ("PHASE 3", ACCENT_CYAN, "Web Dashboard & Multi-User", "React/Next.js frontend for SENTRY. Multi-user access. Visual readiness scoring. Exec committee overview dashboard."),
    ("PHASE 4", ACCENT_RED, "Full Production Deployment", "Discord integration live. All channels connected. Automated alerts. Real-time dashboard. DSTA-ready pitch package.")
]

for i, (phase, phase_color, title, desc) in enumerate(roadmap):
    top = Inches(1.8 + i * 1.1)
    # Phase label
    add_text_box(slide, Inches(0.8), top, Inches(2.2), Inches(0.9),
                 phase, font_size=11, bold=True, color=phase_color)
    # Separator line
    add_shape(slide, Inches(3.0), top, Inches(0.02), Inches(0.8),
              fill_color=phase_color, line_color=phase_color)
    # Title
    add_text_box(slide, Inches(3.2), top, Inches(3.5), Inches(0.4),
                 title, font_size=14, bold=True, color=TEXT_PRIMARY)
    # Description
    add_text_box(slide, Inches(3.2), Inches(top.inches + 0.35), Inches(9.5), Inches(0.5),
                 desc, font_size=11, color=TEXT_SECONDARY)

slide_number(slide, 9)


# ─── SLIDE 10: DSTA ALIGNMENT ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
             "STAKEHOLDER ALIGNMENT", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Why This Fits DSTA Projects",
             font_size=36, bold=True, color=WHITE)

# Four cards
dsta_cards = [
    ("🎯  Problem Fit", "DSTA manages large-scale defence tech projects with multiple stakeholders, committees, and supply chains. Information overload is a known pain point."),
    ("🔐  Security-First Design", "No data exfiltration. Source-tiering built in. Confidence markers for every assessment. Hannah is sole principal with override authority."),
    ("📈  Scalable for Committees", "From individual to full executive committee overview. Task delegation, status tracking, and priority escalation built in."),
    ("🤖  AI-Native, Human-Led", "AI does the filtering, sorting, flagging. Humans make the decisions. The Claw makes you faster — it doesn't replace judgment.")
]

for i, (title, desc) in enumerate(dsta_cards):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(2.2 + row * 2.2)
    card = add_shape(slide, left, top, Inches(5.8), Inches(1.9))
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.font.name = 'Calibri'
    p2.space_before = Pt(10)

# Bottom callout box
box = add_shape(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.9),
                fill_color=RGBColor(0x0F, 0x17, 0x2A), line_color=ACCENT_BLUE)
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Bottom line: A defence/security Claw that ingests project chaos, compresses it into actionable intelligence, tracks what matters, flags what is urgent, checks what is missing, and updates the team — all through a single conversational interface."
p.font.size = Pt(14)
p.font.color.rgb = TEXT_PRIMARY
p.font.name = 'Calibri'
p.font.italic = True

slide_number(slide, 10)


# ─── SAVE ───
output_path = os.path.expanduser('~/.openclaw/workspace/Project-Claw-DSTA-Deck.pptx')
prs.save(output_path)
print(f'Slide count: {len(prs.slides)}')
print(f'Saved to: {output_path}')