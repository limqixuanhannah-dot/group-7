#!/usr/bin/env python3
"""Merged SENTRY deck — 8 slides, all capabilities integrated, no names."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x0A, 0x0E, 0x17)
DARK_CARD = RGBColor(0x11, 0x18, 0x27)
BORDER = RGBColor(0x1E, 0x29, 0x3B)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
TP = RGBColor(0xE0, 0xE6, 0xF0)
TS = RGBColor(0x94, 0xA3, 0xB8)
TM = RGBColor(0x64, 0x74, 0x8B)
WH = RGBColor(0xFF, 0xFF, 0xFF)


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG


def card(slide, l, t, w, h, fc=DARK_CARD, bc=BORDER):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    s.line.color.rgb = bc; s.line.width = Pt(1)
    return s


def tb(slide, l, t, w, h, text, sz=14, b=False, c=TP, a=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.bold = b
    p.font.color.rgb = c; p.font.name = 'Calibri'; p.alignment = a
    return bx.text_frame


def ap(tf, text, sz=14, b=False, c=TS, sb=4):
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.bold = b
    p.font.color.rgb = c; p.font.name = 'Calibri'; p.space_before = Pt(sb)
    return p


def sn(slide, num):
    tb(slide, Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4),
       f"{num:02d}", sz=11, c=TM, a=PP_ALIGN.RIGHT)


# ─── SLIDE 1: TITLE ───
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
   "◆  SENTRY", sz=18, b=True, c=BLUE)
tb(sl, Inches(0.8), Inches(0.9), Inches(5), Inches(0.4),
   "DEFENCE & SECURITY DECISION SUPPORT", sz=10, c=TM)
tb(sl, Inches(0.8), Inches(2.2), Inches(11), Inches(1.2),
   "SENTRY — Crisis Information Triage\n& Ops-Readiness",
   sz=48, b=True, c=WH)
tb(sl, Inches(0.8), Inches(3.8), Inches(10), Inches(0.8),
   "An AI-powered ops assistant that ingests project chatter, triages information,\ntracks tasks, monitors threats, checks logistics, and surfaces what matters.",
   sz=18, c=TS)
tb(sl, Inches(0.8), Inches(4.8), Inches(8), Inches(0.4),
   "9 CAPABILITY DOMAINS  |  12+ OPERATIONAL SYSTEMS  |  HOURLY THREAT MONITORING",
   sz=11, b=True, c=CYAN)
sn(sl, 1)


# ─── SLIDE 2: PROBLEM & SOLUTION ───
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4),
   "PROBLEM & SOLUTION", sz=11, b=True, c=BLUE)
tb(sl, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
   "The Problem It Solves",
   sz=36, b=True, c=WH)
tb(sl, Inches(0.8), Inches(1.7), Inches(11), Inches(0.6),
   "Large projects generate massive cross-channel information flow. Teams and committees struggle to maintain situational awareness.",
   sz=15, c=TS)

# Problem box
bx = card(sl, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.8),
          fc=RGBColor(0x0F, 0x17, 0x2A), bc=BLUE)
tf = bx.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = '"How to make updates, new information, and tasks easier to overview for members and committees of large projects?"'
tf.paragraphs[0].font.size = Pt(15); tf.paragraphs[0].font.color.rgb = TP
tf.paragraphs[0].font.italic = True; tf.paragraphs[0].font.name = 'Calibri'

# Three solution pillars
pillars = [
    ("✂️  Summarise & Filter", "Compress noisy information streams\ninto actionable briefs"),
    ("🚩  Flag Priorities", "Identify urgency, priority, risk —\nCRITICAL → HIGH → MEDIUM → LOW"),
    ("📋  Track Everything", "Tasks, logistics, compliance,\nassets, decisions, threats")
]
for i, (t, d) in enumerate(pillars):
    bx = card(sl, Inches(0.8 + i * 4.1), Inches(3.5), Inches(3.7), Inches(2.0))
    tf = bx.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = t; tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = TP
    tf.paragraphs[0].font.name = 'Calibri'
    ap(tf, d, sz=12, c=TS, sb=8)

# Key metrics
metrics = [
    ("24/7", "Ops Availability"), ("9", "Capability Domains"),
    ("12+", "Systems Deployed"), ("Hourly", "Threat Scanning")]
for i, (n, l) in enumerate(metrics):
    bx = card(sl, Inches(1.5 + i * 2.8), Inches(5.8), Inches(2.4), Inches(1.2))
    tf = bx.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = n; tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER; tf.paragraphs[0].font.name = 'Calibri'
    ap(tf, l, sz=12, c=TM, sb=2); tf.paragraphs[1].alignment = PP_ALIGN.CENTER
sn(sl, 2)


# ─── SLIDE 3: ALL CAPABILITIES ───
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
   "ALL CAPABILITIES", sz=11, b=True, c=BLUE)
tb(sl, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
   "9 Capability Domains — Fully Integrated",
   sz=34, b=True, c=WH)

all_caps = [
    (BLUE, "🛡️  Cyber Hygiene\n& Breach Monitoring", "HIBP integration · Watchlist\n· Real-time alerts"),
    (BLUE, "🚨  Crisis Triage", "Incident playbooks\n· KNOW-DECIDE-ACT"),
    (BLUE, "🔎  Misinformation\n& Fact Check", "Source tiering\n· Confidence markers"),
    (BLUE, "🌐  OSINT & Geopolitical\nAnalysis", "Situation briefs\n· Threat actor profiles"),
    (BLUE, "⚡  Ops Readiness\n& Task Tracking", "Preparedness scoring\n· Priority triage"),
    (CYAN, "📜  Security &\nCompliance", "Classification tagging\n· Deadline tracking"),
    (CYAN, "🚨  Incident Response\nDrills & Templates", "Tabletop scheduling\n· Crisis comms"),
    (CYAN, "🧠  Project\nIntelligence", "Risk scoring · Threat\nbriefings · Vuln triage"),
    (CYAN, "📦  Asset &\nGovernance", "Asset registry · Decision\nlog · Action items"),
]

for i, (clr, t, d) in enumerate(all_caps):
    row = i // 5; col = i % 5 if row == 0 else i - 5
    if row == 0:
        l = Inches(0.8 + col * 2.5); t2 = Inches(2.2)
    else:
        l = Inches(0.8 + col * 2.5 + 1.25); t2 = Inches(4.2)
    bx = card(sl, l, t2, Inches(2.3), Inches(1.8), bc=clr)
    tf = bx.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = t; tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = TP
    tf.paragraphs[0].font.name = 'Calibri'
    ap(tf, d, sz=9, c=TS, sb=4)

# Phase labels
tb(sl, Inches(0.8), Inches(1.8), Inches(12), Inches(0.3),
   "CORE (Phase 1)                                                                                    EXPANDED (Phase 2)",
   sz=9, b=True, c=TM)
sn(sl, 3)


# ─── SLIDE 4: MONITORING & RESEARCH ───
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
   "INTELLIGENCE & MONITORING", sz=11, b=True, c=BLUE)
tb(sl, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
   "Threat Monitoring & Project Research",
   sz=34, b=True, c=WH)
tb(sl, Inches(0.8), Inches(1.7), Inches(11), Inches(0.5),
   "Automated hourly cycles: breach detection + cyber intelligence + project-specific research.",
   sz=15, c=TS)

# Pipeline
stages = ["⏰ Every Hour\nCron Trigger", "🔍 Breach\nCheck (HIBP)", "📰 Intel\nRoundup", "📡 Formatted\nReport"]
for i, s in enumerate(stages):
    bx = card(sl, Inches(0.8 + i * 3.1), Inches(2.3), Inches(2.7), Inches(1.5))
    tf = bx.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = s; tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = CYAN
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Threat intel highlights
tb(sl, Inches(0.8), Inches(4.1), Inches(5), Inches(0.3),
   "KEY BREACHES & THREATS (JUNE 2026)", sz=10, b=True, c=BLUE)

items = [
    ("Canvas/Instructure", "275M users — 3.65 TB — education sector", RED),
    ("Kodak", "2.2M records — ShinyHunters ransomware", RED),
    ("Miasma Worm", "Supply chain — PyPI, npm, GitHub Actions", ORANGE),
    ("Mastra AI npm", "1.1M weekly downloads — malicious code", ORANGE),
    ("FortiSandbox", "3 critical flaws — actively exploited (CISA)", ORANGE),
    ("ShinyHunters Vishing", "100+ orgs — MFA bypass", ORANGE),
]
for i, (n, d, c) in enumerate(items):
    row = i // 3; col = i % 3
    l2 = Inches(0.8 + col * 4.1); t2 = Inches(4.5 + row * 0.7)
    bx = card(sl, l2, t2, Inches(3.8), Inches(0.55), bc=c)
    tf = bx.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = f"{n} — {d}"; tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = c
    tf.paragraphs[0].font.name = 'Calibri'

# Research note
bx = card(sl, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.6),
          fc=RGBColor(0x0F, 0x17, 0x2A), bc=CYAN)
tf = bx.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "Active projects researched hourly: health & wellness (fishkeeping) · productivity tools (Workspace) — new findings appended automatically per project."
tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.color.rgb = CYAN
tf.paragraphs[0].font.name = 'Calibri'
sn(sl, 4)


# ─── SLIDE 5: OPERATIONAL STATUS ───
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
   "OPERATIONAL STATUS", sz=11, b=True, c=BLUE)
tb(sl, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
   "Current System Status",
   sz=34, b=True, c=WH)

# Status grid
st = [
    ("✅  Dashboard", "Operational"),
    ("✅  Heartbeat System", "Online"),
    ("✅  SENTRY Framework", "Deployed"),
    ("✅  Breach Monitoring", "Ready"),
    ("✅  Project Research", "Active"),
    ("⚠️  Discord Delivery", "Needs Token Renewal"),
    ("⚠️  Excel / Sheet Link", "Pending"),
    ("⚠️  Monitoring Targets", "Not Configured"),
]
for i, (lbl, s) in enumerate(st):
    row = i // 4; col = i % 4
    l2 = Inches(0.8 + col * 3.1); t2 = Inches(2.2 + row * 1.3)
    bx = card(sl, l2, t2, Inches(2.8), Inches(1.0))
    tf = bx.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = lbl; tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = TP
    tf.paragraphs[0].font.name = 'Calibri'
    ap(tf, s, sz=10, c=GREEN if "Active" in s or "Online" in s or "Deployed" in s or "Ready" in s or "Operational" in s else ORANGE, sb=4)

# Todo
tb(sl, Inches(0.8), Inches(5.0), Inches(5), Inches(0.3),
   "PENDING ACTIONS", sz=10, b=True, c=BLUE)
todos = [
    ("CRITICAL", "Configure monitoring targets (domains/emails)"),
    ("CRITICAL", "Renew Discord bot token"),
    ("HIGH", "Link project spreadsheet for read/write"),
    ("MEDIUM", "Set up cron for automated reporting"),
]
for i, (sev, t) in enumerate(todos):
    clr = RED if sev == "CRITICAL" else (ORANGE if sev == "HIGH" else CYAN)
    bx = card(sl, Inches(0.8 + (i % 2) * 6.0), Inches(5.4 + (i // 2) * 0.6), Inches(5.5), Inches(0.45), bc=clr)
    tf = bx.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = f"{sev}  —  {t}"; tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = clr
    tf.paragraphs[0].font.name = 'Calibri'
sn(sl, 5)


# ─── SLIDE 6: TECHNICAL ARCHITECTURE ───
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
   "ARCHITECTURE", sz=11, b=True, c=BLUE)
tb(sl, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
   "How It Works",
   sz=34, b=True, c=WH)

arch = [
    ("🧠  AI Core", "DeepSeek V4 reasoning · Sub-agent orchestration\nMemory continuity across sessions · Tool-calling"),
    ("🔧  Tool Layer", "File I/O for spreadsheets · Web search/fetch\nCron scheduling · Parallel session delegation"),
    ("📡  Communications", "WebChat direct · Discord (pending token)\nWeb dashboard · Multi-surface routing"),
    ("🗄️  Data Layer", "Workspace storage · Daily + curated memory\nCSV/Sheet tracking · SENTRY doctrine file"),
]
for i, (t, d) in enumerate(arch):
    row = i // 2; col = i % 2
    l2 = Inches(0.8 + col * 6.2); t2 = Inches(2.0 + row * 2.4)
    bx = card(sl, l2, t2, Inches(5.8), Inches(2.0))
    tf = bx.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = t; tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = TP
    tf.paragraphs[0].font.name = 'Calibri'
    ap(tf, d, sz=12, c=TS, sb=8)

# Tech status
for i, (icon, label) in enumerate([("✓", "Multi-Session"), ("✓", "File I/O"), ("✓", "Web Search"), ("⚠", "Discord"), ("✓", "Cron")]):
    bx = card(sl, Inches(0.8 + i * 2.5), Inches(6.3), Inches(2.2), Inches(0.7))
    tf = bx.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = f"{icon}  {label}"; tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = GREEN if icon == "✓" else ORANGE
    tf.paragraphs[0].font.name = 'Calibri'; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
sn(sl, 6)


# ─── SLIDE 7: SYSTEMS INVENTORY ───
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
   "SYSTEMS INVENTORY", sz=11, b=True, c=BLUE)
tb(sl, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
   "12 Operational Systems — CLI Accessible",
   sz=34, b=True, c=WH)

groups = [
    ("SECURITY & COMPLIANCE", [
        "Classification Tagging (RESTRICTED / CONFIDENTIAL / SECRET)",
        "Compliance Deadline Tracker (auto-detect overdue)",
        "Vendor Security Posture Tracker",
    ], CYAN),
    ("OPS & INCIDENT RESPONSE", [
        "IR Drill Scheduling & Scoring (tabletop / full-scale)",
        "Crisis Communications Templates (breach/outage/disruption)",
    ], ORANGE),
    ("PROJECT INTELLIGENCE", [
        "Risk Scoring Dashboard (4 dimensions, composite score)",
        "Threat Landscape Briefings (per-project correlation)",
        "Vulnerability Disclosure Triage (PENDING → PATCHED)",
    ], RED),
    ("LOGISTICS & GOVERNANCE", [
        "Hardware / Software Asset Registry (expiry, ownership)",
        "Secure File Transfer Log (CREATE / ACCESS / SHARE)",
        "Action Item Tracker + Decision Log (timestamped audit trail)",
    ], GREEN),
]

y = 1.8
for section, items, clr in groups:
    tb(sl, Inches(0.8), Inches(y), Inches(3), Inches(0.3),
       section, sz=9, b=True, c=clr)
    for j, item in enumerate(items):
        tb(sl, Inches(1.2), Inches(y + 0.3 + j * 0.35), Inches(11.5), Inches(0.3),
           f"→  {item}", sz=10, c=TS)
    y += 1.5

bx = card(sl, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.4),
          fc=RGBColor(0x0F, 0x17, 0x2A), bc=BLUE)
tf = bx.text_frame
tf.paragraphs[0].text = "All systems accessible via CLI scripts in scripts/ directory · Cron-ready for automated execution"
tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.color.rgb = BLUE
tf.paragraphs[0].font.name = 'Calibri'; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
sn(sl, 7)


# ─── SLIDE 8: WHY THIS MATTERS ───
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4),
   "PURPOSE & ALIGNMENT", sz=11, b=True, c=BLUE)
tb(sl, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
   "Why This System Exists",
   sz=34, b=True, c=WH)

bx = card(sl, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.7),
          fc=RGBColor(0x0F, 0x17, 0x2A), bc=BLUE)
tf = bx.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = "An integrated ops assistant that ingests project chaos, compresses it into actionable intelligence, tracks what matters, flags what is urgent, monitors cyber threats, and updates the team — through a single conversational interface."
tf.paragraphs[0].font.size = Pt(14); tf.paragraphs[0].font.color.rgb = TP
tf.paragraphs[0].font.italic = True; tf.paragraphs[0].font.name = 'Calibri'

# Purpose pillars
purposes = [
    ("🎯  For Project Teams", [
        "Real-time visibility into status and bottlenecks",
        "No more digging through chat for critical updates",
        "Know what is urgent and what can wait",
    ]),
    ("🏛️  For Committees", [
        "Compressed, accurate briefs over noisy dumps",
        "Confidence-marked assessments, not guessing",
        "Clear escalation path for critical items",
    ]),
    ("🛡️  For Security Ops", [
        "Automated breach detection and cyber hygiene",
        "Source-trusted fact-checking for intelligence",
        "OPSEC-controlled — no data exfiltration",
    ]),
]
for i, (t, items) in enumerate(purposes):
    bx = card(sl, Inches(0.8 + i * 4.1), Inches(2.7), Inches(3.7), Inches(3.0))
    tf = bx.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = t; tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = TP
    tf.paragraphs[0].font.name = 'Calibri'
    for item in items:
        ap(tf, f"→  {item}", sz=11, c=TS, sb=6)

# Outcomes
tb(sl, Inches(0.8), Inches(6.0), Inches(12), Inches(0.3),
   "Three core outcomes:  FASTER DECISIONS  ·  FEWER SURPRISES  ·  BETTER OVERSIGHT",
   sz=13, b=True, c=CYAN, a=PP_ALIGN.CENTER)

# Alignment
bx = card(sl, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.6),
          fc=RGBColor(0x0F, 0x17, 0x2A), bc=BLUE)
tf = bx.text_frame
tf.paragraphs[0].text = "Security-first design · AI filters, humans decide · Scalable from teams to executive committees"
tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.color.rgb = CYAN
tf.paragraphs[0].font.name = 'Calibri'; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
sn(sl, 8)


# ─── SAVE ───
out = os.path.expanduser('~/.openclaw/workspace/SENTRY-Integrated-Deck.pptx')
prs.save(out)
print(f'Slides: {len(prs.slides)}')
for i, sl in enumerate(prs.slides):
    print(f'  Slide {i+1}: {len(list(sl.shapes))} shapes')
print(f'Saved: {out}')
