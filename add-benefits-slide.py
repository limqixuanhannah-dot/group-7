#!/usr/bin/env python3
"""Add an 11th slide: Purpose & Benefits of The Claw, then reorder to position 6."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

path = os.path.expanduser('~/.openclaw/workspace/Project-Claw-DSTA-Deck.pptx')
prs = Presentation(path)

# Colors
DARK_BG = RGBColor(0x0A, 0x0E, 0x17)
DARK_CARD = RGBColor(0x11, 0x18, 0x27)
BORDER = RGBColor(0x1E, 0x29, 0x3B)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)
TEXT_PRIMARY = RGBColor(0xE0, 0xE6, 0xF0)
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=DARK_CARD, line_color=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def slide_number(slide, num):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4),
                 f"{num:02d}", font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)


# ─── CREATE NEW SLIDE: PURPOSE & BENEFITS ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
             "PURPOSE & BENEFITS", font_size=11, bold=True, color=ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Why We Need This Claw",
             font_size=36, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
             "The Claw exists to solve a specific operational problem — and delivers measurable benefits to every stakeholder in the project ecosystem.",
             font_size=15, color=TEXT_SECONDARY)

benefits = [
    (
        "🎯",
        "For Project Leads",
        [
            "Real-time visibility into project status, task progress, and bottlenecks",
            "",
            "No more digging through chat logs to find the one critical update",
            "",
            "Know exactly what's urgent and what can wait"
        ]
    ),
    (
        "🏛️",
        "For Executive Committees",
        [
            "Compressed, accurate briefs instead of noisy information dumps",
            "",
            "Confidence-marked assessments (not AI guessing as fact)",
            "",
            "Escalation path for truly critical items"
        ]
    ),
    (
        "🛡️",
        "For Defence & Security",
        [
            "Breach monitoring and cyber hygiene integrated into daily ops",
            "",
            "Misinformation / fact-check capability for decision confidence",
            "",
            "OPSEC-aware — no data exfiltration, principal-controlled access"
        ]
    )
]

for i, (icon, title, lines) in enumerate(benefits):
    left = Inches(0.8 + i * 4.1)
    card = add_shape(slide, left, Inches(2.6), Inches(3.7), Inches(3.6))
    tf = card.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = f"{icon}  {title}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.font.name = 'Calibri'

    for line in lines:
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_SECONDARY
        p2.font.name = 'Calibri'
        p2.space_before = Pt(2)

add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11), Inches(0.3),
             "Three core outcomes:",
             font_size=12, bold=True, color=TEXT_MUTED)

pillars = [
    ("⚡  FASTER DECISIONS", "Information compressed into actionable briefs"),
    ("🎯  FEWER SURPRISES", "Urgent items surfaced before they become crises"),
    ("📊  BETTER OVERSIGHT", "One interface for status, tasks, logistics, alerts")
]

for i, (title, desc) in enumerate(pillars):
    left = Inches(2.0 + i * 3.5)
    box = add_shape(slide, left, Inches(6.7), Inches(3.0), Inches(0.7),
                    fill_color=DARK_CARD, line_color=BORDER)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = 'Calibri'
    p2.alignment = PP_ALIGN.CENTER

slide_number(slide, 6)


# ─── REORDER: Move new slide (now at end, index 10) to position 6 (index 5) ───
ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
      'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}

pres_elem = prs.element
sldIdLst = pres_elem.find('.//p:sldIdLst', ns)
children = list(sldIdLst)

# Last child is our new slide; remove it
new_slide_elem = children[-1]
sldIdLst.remove(new_slide_elem)

# Insert at position 5 (0-indexed) — after slide 5 (Status), before slide 6 (Chat)
target_index = 5
# Re-fetch children list after removal
current_children = list(sldIdLst)
target = current_children[target_index]
sldIdLst.insert(list(sldIdLst).index(target), new_slide_elem)

print(f"Total slides after reorder: {len(list(sldIdLst))}")

# ─── FIX SLIDE NUMBERS ───
# Slide numbers are in text boxes at Inches(12.2, 7.0) with 2-digit content
for idx, slide in enumerate(prs.slides):
    correct_num = f"{idx + 1:02d}"
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            text = tf.text.strip()
            # Only fix shapes that contain exactly a 2-digit number
            if text.isdigit() and len(text) == 2:
                # Update all paragraphs in this text frame to the correct number
                for p in tf.paragraphs:
                    p.text = correct_num

output_path = os.path.expanduser('~/.openclaw/workspace/Project-Claw-DSTA-Deck.pptx')
prs.save(output_path)
print(f'Saved to: {output_path}')
print(f'Final slide count: {len(prs.slides)}')
for i, slide in enumerate(prs.slides):
    # Get first text to identify the slide
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            t = shape.text_frame.text.strip()[:60]
            texts.append(t)
    # Find the H2-like text (big text)
    big = [t for t in texts if len(t) > 10][:2]
    label = ' | '.join(big) if big else '(no title)'
    print(f'  Slide {i+1}: {label}')
